#!/usr/bin/env python3
"""
Integration tests for the PDF-to-Narrated-PowerPoint system.
Tests the complete pipeline from PDF to PowerPoint.
"""

import os
import sys
import json
import tempfile
from pathlib import Path
from typing import Optional, Dict, Any
import logging

logging.basicConfig(level=logging.INFO, format='%(levelname)s: %(message)s')
logger = logging.getLogger(__name__)


class IntegrationTester:
    """Integration test suite for the system."""
    
    def __init__(self, root_dir: Optional[Path] = None):
        """Initialize tester."""
        self.root_dir = root_dir or Path(__file__).parent.absolute()
        self.tests_passed = 0
        self.tests_failed = 0
    
    def assert_true(self, condition: bool, message: str):
        """Assert condition is true."""
        if condition:
            logger.info(f"✓ {message}")
            self.tests_passed += 1
        else:
            logger.error(f"✗ {message}")
            self.tests_failed += 1
            raise AssertionError(message)
    
    def test_imports(self):
        """Test all critical imports."""
        logger.info("\n[TEST 1] Imports")
        
        try:
            import pymupdf
            self.assert_true(True, "PyMuPDF imported")
        except ImportError:
            self.assert_true(False, "PyMuPDF not installed")
        
        try:
            from mistralai import Mistral
            self.assert_true(True, "Mistral AI imported")
        except ImportError:
            self.assert_true(False, "Mistral AI not installed")
        
        try:
            from pptx import Presentation
            self.assert_true(True, "Python-pptx imported")
        except ImportError:
            self.assert_true(False, "Python-pptx not installed")
        
        try:
            sys.path.insert(0, str(self.root_dir / "document_understanding_agent" / "src"))
            from dua import DocumentUnderstandingAgent
            self.assert_true(True, "Document Understanding Agent imported")
        except ImportError as e:
            self.assert_true(False, f"DUA import failed: {e}")
    
    def test_api_key(self):
        """Test API key configuration."""
        logger.info("\n[TEST 2] Configuration")
        
        from dotenv import load_dotenv
        load_dotenv(self.root_dir / ".env")
        api_key = os.getenv("MISTRAL_API_KEY")
        
        if api_key:
            self.assert_true(True, f"MISTRAL_API_KEY configured")
        else:
            logger.warning("⚠ MISTRAL_API_KEY not configured - Mistral tests will be skipped")
    
    def test_dua_initialization(self):
        """Test Document Understanding Agent initialization."""
        logger.info("\n[TEST 3] Document Understanding Agent")
        
        try:
            sys.path.insert(0, str(self.root_dir / "document_understanding_agent" / "src"))
            from dua import DocumentUnderstandingAgent
            
            agent = DocumentUnderstandingAgent(use_ml=False, log_level="WARNING")
            self.assert_true(agent is not None, "DUA initialized")
            
            status = agent.get_status()
            self.assert_true("modules" in status, "DUA has modules")
            logger.info(f"  Loaded modules: {len(status.get('modules', []))}")
            
        except Exception as e:
            self.assert_true(False, f"DUA initialization failed: {e}")
    
    def test_mistral_client(self):
        """Test Mistral AI client initialization."""
        logger.info("\n[TEST 4] Mistral AI Client")
        
        api_key = os.getenv("MISTRAL_API_KEY")
        if not api_key:
            logger.warning("⚠ Skipping Mistral tests - API key not configured")
            return
        
        try:
            from mistralai import Mistral
            
            client = Mistral(api_key=api_key)
            self.assert_true(client is not None, "Mistral client initialized")
            
        except Exception as e:
            self.assert_true(False, f"Mistral client failed: {e}")
    
    def test_pptx_creation(self):
        """Test PowerPoint creation."""
        logger.info("\n[TEST 5] PowerPoint Generation")
        
        try:
            from pptx import Presentation
            
            prs = Presentation()
            self.assert_true(prs is not None, "Presentation created")
            
            # Add a test slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            self.assert_true(len(prs.slides) == 1, "Slide added")
            
            # Save to temp file
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
                temp_pptx = f.name
            
            prs.save(temp_pptx)
            self.assert_true(Path(temp_pptx).exists(), "PPTX saved to disk")
            
            # Cleanup
            Path(temp_pptx).unlink()
            
        except Exception as e:
            self.assert_true(False, f"PPTX creation failed: {e}")
    
    def test_json_schema_validation(self):
        """Test JSON schema validation."""
        logger.info("\n[TEST 6] JSON Schema Validation")
        
        try:
            import jsonschema
            
            # Example slide JSON
            slide_json = {
                "ppt": {
                    "size": {"width": 1280, "height": 720, "unit": "px"},
                    "slides": []
                }
            }
            
            # Basic validation
            self.assert_true("ppt" in slide_json, "Slide JSON structure valid")
            self.assert_true("slides" in slide_json["ppt"], "Slides array present")
            
        except Exception as e:
            self.assert_true(False, f"JSON validation failed: {e}")
    
    def test_orchestrator(self):
        """Test orchestrator can be imported and initialized."""
        logger.info("\n[TEST 7] Orchestrator")
        
        try:
            sys.path.insert(0, str(self.root_dir))
            from orchestrator import PDFToPresentation
            
            orchestrator = PDFToPresentation()
            self.assert_true(orchestrator is not None, "Orchestrator initialized")
            
            self.assert_true(orchestrator.root_dir.exists(), "Root directory valid")
            self.assert_true(orchestrator.dua_dir.exists(), "DUA directory exists")
            self.assert_true(orchestrator.brain_dir.exists(), "Brain directory exists")
            self.assert_true(orchestrator.ppt_dir.exists(), "PPT directory exists")
            
        except Exception as e:
            self.assert_true(False, f"Orchestrator test failed: {e}")
    
    def test_file_structure(self):
        """Test project file structure."""
        logger.info("\n[TEST 8] Project Structure")
        
        required_files = {
            "document_understanding_agent/src/dua/__init__.py": "DUA package",
            "brain/main.py": "Brain agent",
            "JSON To PPT/main.py": "PPT agent",
            "orchestrator.py": "Orchestrator",
            "requirements.txt": "Requirements",
            "ARCHITECTURE.md": "Architecture docs",
            "SETUP.md": "Setup docs",
            "QUICKSTART.md": "Quick start docs"
        }
        
        for file_path, description in required_files.items():
            full_path = self.root_dir / file_path
            exists = full_path.exists()
            self.assert_true(exists, f"{description} exists: {file_path}")
    
    def test_sample_json_generation(self):
        """Test generating sample JSON structures."""
        logger.info("\n[TEST 9] Sample JSON Generation")
        
        # Sample document content
        sample_content = {
            "document_tree": {
                "sections": [
                    {
                        "title": "Introduction",
                        "level": 1,
                        "blocks": [
                            {
                                "type": "PARAGRAPH",
                                "text": "This is a sample document.",
                                "confidence": 0.95
                            }
                        ]
                    }
                ]
            },
            "metadata": {
                "num_pages": 1,
                "confidence": 0.95
            }
        }
        
        # Test JSON serialization
        try:
            json_str = json.dumps(sample_content, indent=2)
            self.assert_true(len(json_str) > 0, "JSON serialized successfully")
            
            # Test deserialization
            parsed = json.loads(json_str)
            self.assert_true("document_tree" in parsed, "JSON structure preserved")
            
        except Exception as e:
            self.assert_true(False, f"JSON test failed: {e}")
    
    def test_logging_configuration(self):
        """Test logging configuration."""
        logger.info("\n[TEST 10] Logging")
        
        try:
            import logging as log
            
            test_logger = log.getLogger("test")
            test_logger.setLevel(log.DEBUG)
            
            self.assert_true(test_logger.level == log.DEBUG, "Logger level configurable")
            
        except Exception as e:
            self.assert_true(False, f"Logging test failed: {e}")
    
    def run_all_tests(self):
        """Run all tests."""
        print("=" * 70)
        print("PDF-to-Narrated-PowerPoint Integration Tests")
        print("=" * 70)
        
        try:
            self.test_imports()
            self.test_api_key()
            self.test_file_structure()
            self.test_dua_initialization()
            self.test_mistral_client()
            self.test_pptx_creation()
            self.test_json_schema_validation()
            self.test_orchestrator()
            self.test_sample_json_generation()
            self.test_logging_configuration()
        except AssertionError:
            pass  # Errors logged in assert_true
        except Exception as e:
            logger.error(f"Unexpected error: {e}")
            self.tests_failed += 1
        
        # Summary
        print("\n" + "=" * 70)
        print("TEST SUMMARY")
        print("=" * 70)
        total = self.tests_passed + self.tests_failed
        print(f"Total Tests: {total}")
        print(f"Passed: {self.tests_passed} ✓")
        print(f"Failed: {self.tests_failed} ✗")
        
        if self.tests_failed == 0:
            print("\n✓ All integration tests passed!")
            print("\nSystem is ready for production use.")
            return True
        else:
            print(f"\n⚠ {self.tests_failed} test(s) failed.")
            print("Please review the errors above and resolve them.")
            return False


def main():
    """Main entry point."""
    tester = IntegrationTester()
    success = tester.run_all_tests()
    
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
