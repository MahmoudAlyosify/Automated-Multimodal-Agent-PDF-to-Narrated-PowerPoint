"""pptx_builder.py

A single-file, trimmed extractor of the project's core JSON->PPT logic.

Usage:
  from pptx_builder import build, validate, save_to_temp
  prs, count = build(meta_dict)
  path = save_to_temp(prs)

This file intentionally contains a compact, accurate subset of the original
`main.py` logic: layout/size handling, text/image/shape/table/chart rendering,
image sources (base64,url,file), and basic style helpers.
"""
import base64
import io
import json
import os
import re
import uuid
import logging
from typing import Any, Dict, Optional, Tuple, List, Set

try:
    import requests
except Exception:
    requests = None

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.dml import MSO_LINE
except Exception:
    Presentation = None

# Limits and config
MAX_SLIDES = int(os.environ.get("JSON2PPT_MAX_SLIDES", "200"))
MAX_ELEMENTS_PER_SLIDE = int(os.environ.get("JSON2PPT_MAX_ELEMENTS", "600"))
MAX_REMOTE_IMAGE_BYTES = int(os.environ.get("JSON2PPT_MAX_IMG_BYTES", str(5 * 1024 * 1024)))
ALLOW_REMOTE_IMAGES = os.environ.get("JSON2PPT_ALLOW_REMOTE", "1") == "1"
ALLOW_FILE_IMAGES = os.environ.get("JSON2PPT_ALLOW_FILE", "1") == "1"
ASSET_ROOT = os.environ.get("JSON2PPT_ASSET_ROOT")
ALLOWED_REMOTE_DOMAINS: Optional[Set[str]] = None
if os.environ.get("JSON2PPT_REMOTE_DOMAINS"):
    ALLOWED_REMOTE_DOMAINS = {d.strip().lower() for d in os.environ["JSON2PPT_REMOTE_DOMAINS"].split(",") if d.strip()}

_IMAGE_CACHE: Dict[str, bytes] = {}
_THEME: Dict[str, Any] = {}


def hex_to_rgb(color: str):
    if not color:
        return (0, 0, 0)
    c = color.strip().lstrip('#')
    if len(c) == 3:
        c = ''.join(ch * 2 for ch in c)
    if len(c) >= 6:
        return (int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
    return (0, 0, 0)


def unit_to_emu(value: float, total: int, unit: str) -> int:
    if unit == "percent":
        px = total * value / 100.0
    else:
        px = value
    return int(px * 9525)


def resolve_box(box: Dict[str, Any], slide_w: int, slide_h: int, default_unit: str):
    unit = box.get("unit", default_unit)
    x = unit_to_emu(box.get("x", 0), slide_w, unit)
    y = unit_to_emu(box.get("y", 0), slide_h, unit)
    w = unit_to_emu(box.get("w", slide_w), slide_w, unit)
    h = unit_to_emu(box.get("h", slide_h), slide_h, unit)
    return x, y, w, h


def _domain_allowed(url: str) -> bool:
    if not ALLOWED_REMOTE_DOMAINS:
        return True
    try:
        m = re.match(r'^https?://([^/]+)', url, re.IGNORECASE)
        if not m:
            return False
        host = m.group(1).lower()
        return any(host == d or host.endswith('.' + d) for d in ALLOWED_REMOTE_DOMAINS)
    except Exception:
        return False


def _check_file_path(path: str) -> bool:
    if not ALLOW_FILE_IMAGES:
        return False
    if ASSET_ROOT:
        try:
            real_root = os.path.realpath(ASSET_ROOT)
            real_path = os.path.realpath(path)
            return real_path.startswith(real_root)
        except Exception:
            return False
    return True


def get_image_bytes(source: str, logger=None):
    if not source:
        return None
    cached = _IMAGE_CACHE.get(source)
    if cached is not None:
        return cached
    try:
        if source.startswith("base64:"):
            b64 = source[len("base64:"):]
            missing = len(b64) % 4
            if missing:
                b64 += '=' * (4 - missing)
            content = base64.b64decode(b64)
            _IMAGE_CACHE[source] = content
            return content
        elif source.startswith("url:"):
            if not ALLOW_REMOTE_IMAGES:
                logger and logger.warning("remote image disabled by config")
                return None
            if requests is None:
                logger and logger.warning("requests not installed")
                return None
            url = source[len("url:"):]
            if not url.lower().startswith(("http://", "https://")):
                logger and logger.warning("unsupported url scheme")
                return None
            if not _domain_allowed(url):
                logger and logger.warning(f"domain not allowed: {url}")
                return None
            resp = requests.get(url, timeout=(3, 7), stream=True)
            resp.raise_for_status()
            data = b''
            total = 0
            for chunk in resp.iter_content(65536):
                if not chunk:
                    continue
                data += chunk
                total += len(chunk)
                if total > MAX_REMOTE_IMAGE_BYTES:
                    logger and logger.warning("remote image exceeds size limit")
                    return None
            _IMAGE_CACHE[source] = data
            return data
        elif source.startswith("file:"):
            path = source[len("file:"):]
            if not _check_file_path(path):
                logger and logger.warning("file path not allowed")
                return None
            with open(path, "rb") as f:
                content = f.read()
                if len(content) > MAX_REMOTE_IMAGE_BYTES:
                    logger and logger.warning("local image exceeds size limit")
                    return None
                _IMAGE_CACHE[source] = content
                return content
    except Exception as e:
        logger and logger.warning(f"get_image_bytes failed: {e}")
    return None


def ppt_color(fill_obj, color_hex: str):
    if not color_hex:
        return
    r, g, b = hex_to_rgb(color_hex)[:3]
    fill_obj.solid()
    fill_obj.fore_color.rgb = RGBColor(r, g, b)


def map_align(a: str):
    if a == "center":
        return PP_ALIGN.CENTER
    if a == "right":
        return PP_ALIGN.RIGHT
    if a == "justify":
        return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT


def map_valign(v: str):
    if v == "middle":
        return MSO_ANCHOR.MIDDLE
    if v == "bottom":
        return MSO_ANCHOR.BOTTOM
    return MSO_ANCHOR.TOP


def merge_styles(base: Dict[str, Any], override: Dict[str, Any]) -> Dict[str, Any]:
    merged = dict(base or {})
    merged.update(override or {})
    return merged


def apply_run_style(run, style: Dict[str, Any]):
    if not style:
        return
    if "fontSize" in style:
        run.font.size = Pt(style["fontSize"])
    if style.get("bold") is not None:
        run.font.bold = bool(style.get("bold"))
    if style.get("italic") is not None:
        run.font.italic = bool(style.get("italic"))
    if style.get("color"):
        r, g, b = hex_to_rgb(style["color"])[:3]
        run.font.color.rgb = RGBColor(r, g, b)


def apply_paragraph_style(paragraph, style: Dict[str, Any]):
    if not style:
        return
    if style.get("align"):
        paragraph.alignment = map_align(style.get("align"))
    if paragraph.runs:
        run = paragraph.runs[0]
    else:
        run = paragraph.add_run()
    apply_run_style(run, style)


def apply_cell_style(cell, style: Dict[str, Any]):
    if not style:
        return
    if style.get("fill"):
        ppt_color(cell.fill, style["fill"])
    if style.get("valign"):
        cell.text_frame.vertical_anchor = map_valign(style.get("valign"))
    for paragraph in cell.text_frame.paragraphs:
        apply_paragraph_style(paragraph, style)


def pick_shape(shape_type: str):
    mapping = {
        "rect": MSO_SHAPE.RECTANGLE,
        "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL,
        "circle": MSO_SHAPE.OVAL,
    }
    return mapping.get(shape_type, MSO_SHAPE.RECTANGLE)


def map_chart_type(t: str):
    mapping = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "scatter": XL_CHART_TYPE.XY_SCATTER_LINES,
        "bubble": XL_CHART_TYPE.BUBBLE,
    }
    return mapping.get(t, XL_CHART_TYPE.COLUMN_CLUSTERED)


def apply_base_element_properties(shape, elem: Dict[str, Any]):
    if not shape or not elem:
        return
    if "rotation" in elem:
        try:
            shape.rotation = elem["rotation"]
        except Exception:
            pass
    if "opacity" in elem:
        try:
            if hasattr(shape, 'fill'):
                shape.fill.transparency = 1.0 - float(elem["opacity"]) if elem["opacity"] is not None else 1.0
        except Exception:
            pass


def add_background(slide, bg_cfg: Dict[str, Any], slide_w_emu: int, slide_h_emu: int, logger=None):
    if not bg_cfg:
        return
    color = bg_cfg.get("color")
    if color:
        ppt_color(slide.background.fill, color)
    img = bg_cfg.get("image")
    if img and isinstance(img, dict):
        content = get_image_bytes(img.get("src"), logger)
        if content:
            pic = slide.shapes.add_picture(io.BytesIO(content), 0, 0, width=slide_w_emu, height=slide_h_emu)
            try:
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.insert(2, pic._element)
            except Exception:
                pass


def add_text(slide, elem, slide_w, slide_h, default_unit):
    box = elem.get("box", {})
    x, y, w, h = resolve_box(box, slide_w, slide_h, default_unit)
    shape = slide.shapes.add_textbox(x, y, w, h)
    apply_base_element_properties(shape, elem)
    if elem.get("fill"):
        ppt_color(shape.fill, elem["fill"])
    tf = shape.text_frame
    tf.clear()
    style = elem.get("style", {})
    paragraphs = elem.get("paragraphs")
    if paragraphs:
        for idx, para_cfg in enumerate(paragraphs):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            text = para_cfg.get("text", "")
            p.text = text
            apply_paragraph_style(p, merge_styles(style, para_cfg.get("style")))
    else:
        p = tf.paragraphs[0]
        p.text = elem.get("text", "")
        apply_paragraph_style(p, style)


def add_image(slide, elem, slide_w, slide_h, default_unit, logger=None):
    box = elem.get("box", {})
    x, y, w, h = resolve_box(box, slide_w, slide_h, default_unit)
    content = get_image_bytes(elem.get("source") or elem.get("src"), logger)
    if not content:
        return
    pic = slide.shapes.add_picture(io.BytesIO(content), x, y, width=w, height=h)
    apply_base_element_properties(pic, elem)


def add_shape(slide, elem, slide_w, slide_h, default_unit):
    box = elem.get("box", {})
    x, y, w, h = resolve_box(box, slide_w, slide_h, default_unit)
    shape_type = elem.get("shapeType", "rect")
    if shape_type == "circle":
        size = min(w, h)
        w = h = size
    shape = slide.shapes.add_shape(pick_shape(shape_type), x, y, w, h)
    apply_base_element_properties(shape, elem)
    if elem.get("fill"):
        ppt_color(shape.fill, elem["fill"])
    if elem.get("border"):
        try:
            line = shape.line
            line.color.rgb = RGBColor(*hex_to_rgb(elem["border"].get("color", "#000000"))[:3])
        except Exception:
            pass


def add_table(slide, elem, slide_w, slide_h, default_unit):
    box = elem.get("box", {})
    x, y, w, h = resolve_box(box, slide_w, slide_h, default_unit)
    table_cfg = elem.get("table", {})
    header = table_cfg.get("header") or []
    rows = table_cfg.get("rows") or []
    column_count = len(header) if header else (len(rows[0]) if rows else 0)
    if column_count == 0:
        return
    total_rows = len(rows) + (1 if header else 0)
    shape = slide.shapes.add_table(total_rows, column_count, x, y, w, h)
    table = shape.table
    def populate_row(row_idx, values):
        for col_idx in range(column_count):
            cell = table.cell(row_idx, col_idx)
            text = values[col_idx] if col_idx < len(values) else ""
            cell.text = "" if text is None else str(text)
    r = 0
    if header:
        populate_row(r, header); r += 1
    for rv in rows:
        populate_row(r, rv); r += 1


def add_line(slide, elem, slide_w, slide_h, default_unit):
    pts = elem.get("points", [])
    if len(pts) < 2:
        return
    p1 = pts[0]; p2 = pts[1]
    x1 = unit_to_emu(p1.get("x", 0), slide_w, "px")
    y1 = unit_to_emu(p1.get("y", 0), slide_h, "px")
    x2 = unit_to_emu(p2.get("x", 100), slide_w, "px")
    y2 = unit_to_emu(p2.get("y", 100), slide_h, "px")
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    try:
        if elem.get("stroke"):
            connector.line.color.rgb = RGBColor(*hex_to_rgb(elem.get("stroke"))[:3])
    except Exception:
        pass


def add_chart(slide, elem, slide_w, slide_h, default_unit, logger=None):
    box = elem.get("box", {})
    x, y, w, h = resolve_box(box, slide_w, slide_h, default_unit)
    chart_type = map_chart_type(elem.get("chartType", "bar"))
    data_cfg = elem.get("data", {}) or {}
    raw_categories = data_cfg.get("categories") or []
    series_cfg = data_cfg.get("series") or []
    try:
        chart_name = getattr(chart_type, 'name', str(chart_type))
        is_xy = 'SCATTER' in chart_name or 'XY_SCATTER' in chart_name
        is_bubble = 'BUBBLE' in chart_name
        if is_bubble:
            bdata = BubbleChartData()
            for s in series_cfg or [{"name": "Series", "values": [(0, 0, 1)]}]:
                serie = bdata.add_series(s.get("name", "Series"))
                for v in s.get("values", []):
                    if isinstance(v, (list, tuple)) and len(v) >= 3:
                        try:
                            serie.add_data_point(float(v[0]), float(v[1]), float(v[2]))
                        except Exception:
                            pass
            data_obj = bdata
        elif is_xy:
            xydata = XyChartData()
            for s in series_cfg or [{"name": "Series", "values": [(0, 0)]}]:
                serie = xydata.add_series(s.get("name", "Series"))
                for pair in s.get("values", []):
                    try:
                        serie.add_data_point(float(pair[0]), float(pair[1]))
                    except Exception:
                        pass
            data_obj = xydata
        else:
            cdata = CategoryChartData()
            if not raw_categories and series_cfg:
                longest = max((len(s.get('values', [])) for s in series_cfg), default=0)
                raw_categories = [f'C{i+1}' for i in range(longest)]
            cdata.categories = [str(c) for c in raw_categories]
            for s in series_cfg or [{"name": "Series", "values": [0]}]:
                values = []
                for v in s.get('values', []):
                    try:
                        values.append(float(v))
                    except Exception:
                        values.append(None)
                cdata.add_series(s.get('name', 'Series'), values)
            data_obj = cdata

        chart_shape = slide.shapes.add_chart(chart_type, x, y, w, h, data_obj)
        return chart_shape
    except Exception as e:
        logger and logger.warning(f"add_chart failed: {e}")
        return None


def add_group(slide, elem, slide_w, slide_h, default_unit, logger=None):
    for sub in elem.get('elements', []) or []:
        t = sub.get('type')
        try:
            if t == 'text':
                add_text(slide, sub, slide_w, slide_h, default_unit)
            elif t == 'image':
                add_image(slide, sub, slide_w, slide_h, default_unit, logger)
            elif t == 'shape':
                add_shape(slide, sub, slide_w, slide_h, default_unit)
            elif t == 'table':
                add_table(slide, sub, slide_w, slide_h, default_unit)
            elif t == 'chart':
                add_chart(slide, sub, slide_w, slide_h, default_unit, logger)
            elif t == 'line':
                add_line(slide, sub, slide_w, slide_h, default_unit)
        except Exception as e:
            logger and logger.warning(f"group sub-element failed: {e}")


def resolve_slide_layout(prs, layout_cfg, fallback):
    if layout_cfg is None:
        return fallback
    if isinstance(layout_cfg, int):
        try:
            return prs.slide_layouts[layout_cfg]
        except Exception:
            return fallback
    if isinstance(layout_cfg, str) and layout_cfg.isdigit():
        return resolve_slide_layout(prs, int(layout_cfg), fallback)
    return fallback


def validate(meta: Dict[str, Any], use_schema: bool = False, logger=None):
    if "ppt" not in meta:
        raise ValueError("missing ppt root")
    slides = meta['ppt'].get('slides')
    if not isinstance(slides, list):
        raise ValueError('ppt.slides must be list')
    if len(slides) > MAX_SLIDES:
        raise ValueError('slides exceed limit')
    for i, s in enumerate(slides):
        elems = s.get('elements', [])
        if not isinstance(elems, list):
            raise ValueError(f'slide[{i}].elements must be list')
        if len(elems) > MAX_ELEMENTS_PER_SLIDE:
            raise ValueError(f'slide[{i}] elements exceed limit')
    return True


def build(meta: Dict[str, Any], logger=None):
    if Presentation is None:
        raise RuntimeError('python-pptx not installed')
    prs = Presentation()
    ppt_cfg = meta['ppt']
    global _THEME
    _THEME = ppt_cfg.get('theme', {})
    size_cfg = ppt_cfg.get('size', {'width': 1280, 'height': 720, 'unit': 'px'})
    sw = size_cfg.get('width', 1280)
    sh = size_cfg.get('height', 720)
    prs.slide_width = unit_to_emu(sw, sw, 'px')
    prs.slide_height = unit_to_emu(sh, sh, 'px')
    default_unit = ppt_cfg.get('defaultUnit', 'px')
    default_layout = resolve_slide_layout(prs, ppt_cfg.get('defaultLayout'), prs.slide_layouts[6])
    slides_cfg = ppt_cfg['slides']
    for s in slides_cfg:
        layout = resolve_slide_layout(prs, s.get('layout'), default_layout)
        slide = prs.slides.add_slide(layout)
        add_background(slide, s.get('background'), prs.slide_width, prs.slide_height, logger)
        elements = s.get('elements', [])
        for elem in sorted(elements, key=lambda e: e.get('zIndex', 0)):
            t = elem.get('type')
            try:
                if t == 'text':
                    add_text(slide, elem, sw, sh, default_unit)
                elif t == 'image':
                    add_image(slide, elem, sw, sh, default_unit, logger)
                elif t == 'shape':
                    add_shape(slide, elem, sw, sh, default_unit)
                elif t == 'table':
                    add_table(slide, elem, sw, sh, default_unit)
                elif t == 'chart':
                    add_chart(slide, elem, sw, sh, default_unit, logger)
                elif t == 'line':
                    add_line(slide, elem, sw, sh, default_unit)
                elif t == 'group':
                    add_group(slide, elem, sw, sh, default_unit, logger)
            except Exception as e:
                logger and logger.warning(f"element failed id={elem.get('id')} err={e}")
    return prs, len(slides_cfg)


def save_to_temp(prs) -> str:
    name = f"ppt_{uuid.uuid4().hex}.pptx"
    path = os.path.abspath(os.path.join("./", name))
    prs.save(path)
    return path


if __name__ == '__main__':
    import sys
    if len(sys.argv) < 2:
        print('Usage: python pptx_builder.py sample.json')
        sys.exit(1)
    with open(sys.argv[1], 'r', encoding='utf-8') as f:
        meta = json.load(f)
    validate(meta)
    prs, cnt = build(meta, logger=logging.getLogger())
    out = save_to_temp(prs)
    print(f'Generated {out} ({cnt} slides)')
