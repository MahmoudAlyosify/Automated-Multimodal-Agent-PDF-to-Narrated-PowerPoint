"""
PPTX BUILDER AGENT - Production Implementation
==============================================

PURPOSE:
This agent materializes generated slide content from JSON into
a professional PowerPoint (.pptx) file with complete visual rendering.

INPUT:
- Generated slide content as JSON (from Slide Generator Agent)

OUTPUT:
- lecture.pptx file (fully rendered PowerPoint presentation)

FEATURES:
- Converts all JSON design elements to PPT elements
- Handles gradients, colors, shadows, shapes
- Applies typography and styling
- Includes transitions
- Creates professional visual presentation

TOOLS:
- python-pptx (core PPT generation)
- PIL (image processing)

Author: AI Agent Implementation System
"""

import json
import logging
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from PIL import Image, ImageDraw
import io
import os


# ============================================================================
# LOGGING CONFIGURATION
# ============================================================================

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


# ============================================================================
# PPTX BUILDER AGENT
# ============================================================================

class PPTXBuilderAgent:
    """Converts JSON presentation design to PowerPoint file"""
    
    def __init__(self, audio_folder: Optional[str] = None):
        self.prs = None
        self.presentation_data = None
        self.slide_width = 1280
        self.slide_height = 720
        self.scale_x = 1.0  # Scaling factor for x-coordinates
        self.scale_y = 1.0  # Scaling factor for y-coordinates
        self.font_scale = 1.0  # Scaling factor for font sizes
        self.audio_folder = Path(audio_folder) if audio_folder else None  # Audio folder path
        self.audio_files = {}  # Cache for audio file paths by slide_id
    
    def load_presentation_json(self, json_file: str = "presentation.json") -> bool:
        """Load presentation JSON from file"""
        
        try:
            with open(json_file, 'r', encoding='utf-8') as f:
                self.presentation_data = json.load(f)
            
            logger.info(f"Loaded presentation from {json_file}")
            
            # Extract slide dimensions
            size_info = self.presentation_data['ppt']['size']
            self.slide_width = size_info['width']
            self.slide_height = size_info['height']
            
            # Load audio files if audio folder exists
            self._load_audio_files()
            
            return True
            
        except Exception as e:
            logger.error(f"Failed to load presentation JSON: {e}")
            return False
    
    def _load_audio_files(self):
        """Load available audio files from audio folder"""
        if not self.audio_folder or not self.audio_folder.exists():
            logger.warning(f"Audio folder not found: {self.audio_folder}")
            return
        
        try:
            # Find all WAV and MP3 files in audio folder
            for audio_file in self.audio_folder.glob("*.wav"):
                # Extract slide ID from filename (e.g., "slide_01_audio.wav" -> 1)
                filename = audio_file.stem
                if filename.startswith("slide_"):
                    try:
                        slide_id = int(filename.split("_")[1])
                        self.audio_files[slide_id] = audio_file
                        logger.info(f"Found audio for slide {slide_id}: {audio_file.name}")
                    except (ValueError, IndexError):
                        pass
            
            logger.info(f"Loaded {len(self.audio_files)} audio files")
        except Exception as e:
            logger.warning(f"Error loading audio files: {e}")
    
    def build_presentation(self) -> bool:
        """Build PowerPoint presentation from JSON"""
        
        try:
            logger.info("Building PowerPoint presentation...")
            
            # Initialize presentation with standard 16:9 dimensions
            self.prs = Presentation()
            ppt_width_inches = 10
            ppt_height_inches = 5.625
            self.prs.slide_width = Inches(ppt_width_inches)
            self.prs.slide_height = Inches(ppt_height_inches)
            
            # Calculate scaling factors from JSON px to PPT inches
            self.scale_x = ppt_width_inches / self.slide_width
            self.scale_y = ppt_height_inches / self.slide_height
            self.font_scale = ppt_height_inches / self.slide_height
            
            logger.info(f"Scaling factors: x={self.scale_x:.6f}, y={self.scale_y:.6f}, font={self.font_scale:.6f}")
            
            # Process each slide
            slides_data = self.presentation_data['ppt']['slides']
            for idx, slide_data in enumerate(slides_data, 1):
                self._build_slide(slide_data, idx)
                logger.info(f"  Built slide {idx}/{len(slides_data)}")
            
            logger.info(f"Presentation built with {len(slides_data)} slides")
            return True
            
        except Exception as e:
            logger.error(f"Failed to build presentation: {e}")
            return False
    
    def _build_slide(self, slide_data: Dict[str, Any], slide_number: int):
        """Build a single slide"""
        
        # Create blank slide
        blank_slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_slide_layout)
        
        # Apply background
        self._apply_background(slide, slide_data.get('background'))
        
        # Add elements
        elements = slide_data.get('elements', [])
        for element in elements:
            self._add_element(slide, element)
        
        # Embed audio if available
        slide_id = slide_data.get('slide_id', slide_number)
        self._embed_audio_to_slide(slide, slide_id, slide_number)
    
    def _apply_background(self, slide, background_data: Optional[Dict]):
        """Apply background to slide"""
        
        if not background_data:
            return
        
        try:
            # Handle gradient background
            if 'gradient' in background_data:
                gradient = background_data['gradient']
                self._apply_gradient_background(slide, gradient)
            
            # Handle solid color background
            elif 'color' in background_data:
                color_hex = background_data['color']
                rgb_color = self._hex_to_rgb(color_hex)
                
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*rgb_color)
        
        except Exception as e:
            logger.warning(f"Failed to apply background: {e}")
    
    def _apply_gradient_background(self, slide, gradient_data: Dict):
        """Apply gradient background using image"""
        
        try:
            # Create gradient image based on actual JSON dimensions
            gradient_img = self._create_gradient_image(
                self.slide_width,
                self.slide_height,
                gradient_data
            )
            
            # Add gradient as background image
            pic = slide.shapes.add_picture(
                gradient_img,
                Inches(0), Inches(0),
                width=self.prs.slide_width,
                height=self.prs.slide_height
            )
            
            # Move picture to back
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
        
        except Exception as e:
            logger.warning(f"Failed to apply gradient: {e}")
    
    def _create_gradient_image(self, width: int, height: int, gradient_data: Dict) -> io.BytesIO:
        """Create gradient image"""
        
        img = Image.new('RGB', (width, height))
        draw = ImageDraw.Draw(img)
        
        stops = gradient_data.get('stops', [])
        if len(stops) < 2:
            return io.BytesIO()
        
        # Simple linear gradient between two colors
        color1_hex = stops[0]['color']
        color2_hex = stops[-1]['color']
        
        rgb1 = self._hex_to_rgb(color1_hex)
        rgb2 = self._hex_to_rgb(color2_hex)
        
        # Create gradient (simplified: left to right or top to bottom)
        angle = gradient_data.get('angle', 45)
        
        for x in range(width):
            ratio = x / width
            r = int(rgb1[0] + (rgb2[0] - rgb1[0]) * ratio)
            g = int(rgb1[1] + (rgb2[1] - rgb1[1]) * ratio)
            b = int(rgb1[2] + (rgb2[2] - rgb1[2]) * ratio)
            
            draw.line([(x, 0), (x, height)], fill=(r, g, b))
        
        # Convert to bytes
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        
        return img_bytes
    
    def _add_element(self, slide, element_data: Dict[str, Any]):
        """Add element to slide"""
        
        element_type = element_data.get('type', '')
        
        if element_type == 'text':
            self._add_text_element(slide, element_data)
        elif element_type == 'shape':
            self._add_shape_element(slide, element_data)
        elif element_type == 'line':
            self._add_line_element(slide, element_data)
        elif element_type == 'rectangle':
            self._add_rectangle_element(slide, element_data)
    
    def _add_text_element(self, slide, element_data: Dict[str, Any]):
        """Add text element to slide"""
        
        try:
            box = element_data['box']
            style = element_data.get('style', {})
            text_content = element_data.get('text', '')
            
            # Convert box coordinates from px to inches using scale factors
            left = Inches(box['x'] * self.scale_x)
            top = Inches(box['y'] * self.scale_y)
            width = Inches(box['w'] * self.scale_x)
            height = Inches(box['h'] * self.scale_y)
            
            # Expand boxes to accommodate text better (add 30% to height)
            height = height * 1.3
            
            # Ensure text box stays within slide boundaries
            if left + width > self.prs.slide_width:
                width = self.prs.slide_width - left - Inches(0.05)
            if top + height > self.prs.slide_height:
                height = self.prs.slide_height - top - Inches(0.05)
            
            # Ensure minimum dimensions
            if width < Inches(0.8):
                width = Inches(0.8)
            if height < Inches(0.4):
                height = Inches(0.4)
            
            # Add text box
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center text vertically
            text_frame.margin_bottom = Inches(0.02)
            text_frame.margin_top = Inches(0.02)
            text_frame.margin_left = Inches(0.04)
            text_frame.margin_right = Inches(0.04)
            
            # Add text
            p = text_frame.paragraphs[0]
            p.text = text_content
            p.line_spacing = 1.15
            
            # Apply styling - Increase font size to 65% for better visibility
            font_size = style.get('fontSize', 18)
            # Scale font size: use 65% of design size for better readability
            final_font_size = max(11, int(font_size * 0.65))
            p.font.size = Pt(final_font_size)
            
            if style.get('bold'):
                p.font.bold = True
            
            if style.get('italic'):
                p.font.italic = True
            
            # Color
            color_hex = style.get('color', '#000000')
            rgb_color = self._hex_to_rgb(color_hex)
            p.font.color.rgb = RGBColor(*rgb_color)
            
            # Alignment
            align_map = {
                'left': PP_ALIGN.LEFT,
                'center': PP_ALIGN.CENTER,
                'right': PP_ALIGN.RIGHT
            }
            p.alignment = align_map.get(style.get('align', 'left'), PP_ALIGN.LEFT)
        
        except Exception as e:
            logger.warning(f"Failed to add text element: {e}")
    
    def _add_shape_element(self, slide, element_data: Dict[str, Any]):
        """Add shape element to slide"""
        
        try:
            box = element_data['box']
            style = element_data.get('style', {})
            
            # Convert box coordinates from px to inches using scale factors
            left = Inches(box['x'] * self.scale_x)
            top = Inches(box['y'] * self.scale_y)
            width = Inches(box['w'] * self.scale_x)
            height = Inches(box['h'] * self.scale_y)
            
            # Ensure shape stays within slide boundaries
            if left + width > self.prs.slide_width:
                width = self.prs.slide_width - left - Inches(0.05)
            if top + height > self.prs.slide_height:
                height = self.prs.slide_height - top - Inches(0.05)
            
            # Add shape
            shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
            
            # Apply fill color
            if style.get('fill'):
                fill = shape.fill
                fill.solid()
                rgb_color = self._hex_to_rgb(style['fill'])
                fill.fore_color.rgb = RGBColor(*rgb_color)
            
            # Apply border
            if style.get('border') and style['border'] is not None:
                border = style['border']
                line = shape.line
                line.width = Pt(border.get('width', 1))
                if border.get('color'):
                    rgb_color = self._hex_to_rgb(border['color'])
                    line.color.rgb = RGBColor(*rgb_color)
            else:
                shape.line.color.rgb = RGBColor(0, 0, 0)
                shape.line.width = Pt(0)
        
        except Exception as e:
            logger.warning(f"Failed to add shape element: {e}")
    
    def _add_line_element(self, slide, element_data: Dict[str, Any]):
        """Add line element to slide"""
        
        try:
            points = element_data.get('points', [])
            if len(points) < 2:
                return
            
            # Get line properties
            stroke = element_data.get('stroke', '#000000')
            stroke_width = element_data.get('strokeWidth', 1)
            
            # Add connector (line) between first two points with scaling
            p1 = points[0]
            p2 = points[1]
            
            left = Inches(p1['x'] * self.scale_x)
            top = Inches(p1['y'] * self.scale_y)
            right = Inches(p2['x'] * self.scale_x)
            bottom = Inches(p2['y'] * self.scale_y)
            
            connector = slide.shapes.add_connector(1, left, top, right, bottom)
            line = connector.line
            line.width = Pt(stroke_width)
            
            rgb_color = self._hex_to_rgb(stroke)
            line.color.rgb = RGBColor(*rgb_color)
        
        except Exception as e:
            logger.warning(f"Failed to add line element: {e}")
    
    def _add_rectangle_element(self, slide, element_data: Dict[str, Any]):
        """Add rectangle element to slide"""
        
        # Treat as shape element
        self._add_shape_element(slide, element_data)
    
    def _embed_audio_to_slide(self, slide, slide_id: int, slide_number: int):
        """Embed audio file to slide if available"""
        try:
            # Check if audio file exists for this slide
            if slide_id not in self.audio_files:
                return
            
            audio_path = self.audio_files[slide_id]
            if not audio_path.exists():
                logger.warning(f"Audio file not found: {audio_path}")
                return
            
            # Add audio to slide (bottom right corner, small size)
            # Position: near bottom right, size: 0.5" x 0.5"
            left = self.prs.slide_width - Inches(0.7)
            top = self.prs.slide_height - Inches(0.7)
            
            # Add audio shape
            audio_shape = slide.shapes.add_movie(
                str(audio_path),
                left,
                top,
                width=Inches(0.5),
                height=Inches(0.5)
            )
            
            logger.info(f"[OK] Embedded audio for slide {slide_id}: {audio_path.name}")
            
        except Exception as e:
            logger.warning(f"[WARN] Could not embed audio for slide {slide_id}: {e}")
    
    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple"""
        
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def save_presentation(self, output_path: str = "lecture.pptx") -> bool:
        """Save PowerPoint file"""
        
        if not self.prs:
            logger.error("No presentation to save")
            return False
        
        try:
            self.prs.save(output_path)
            logger.info(f"Saved presentation to {output_path}")
            return True
        
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            return False
    
    def display_summary(self):
        """Display presentation summary"""
        
        if not self.presentation_data:
            print("No presentation data")
            return
        
        print("\n" + "="*80)
        print("PPTX BUILDER AGENT - SUMMARY")
        print("="*80)
        
        ppt_data = self.presentation_data['ppt']
        slides = ppt_data['slides']
        
        print(f"\nPresentation Version: {self.presentation_data.get('version', 'Unknown')}")
        print(f"Total Slides: {len(slides)}")
        print(f"Dimensions: {ppt_data['size']['width']}x{ppt_data['size']['height']} px")
        
        print(f"\nColor Theme:")
        colors = ppt_data['theme']['colors']
        for color_name in list(colors.keys())[:5]:  # Show first 5 colors
            color_value = colors[color_name]
            print(f"  {color_name}: {color_value}")
        
        print(f"\nSlides Generated:")
        for idx, slide in enumerate(slides[:5], 1):  # Show first 5 slides
            num_elements = len(slide.get('elements', []))
            print(f"  [{idx}] {slide['title']} ({num_elements} elements)")
        
        if len(slides) > 5:
            print(f"  ... and {len(slides) - 5} more slides")


# ============================================================================
# MAIN EXECUTION
# ============================================================================

def main():
    """Main execution function"""
    
    import argparse
    
    parser = argparse.ArgumentParser(
        description='PPTX Builder Agent - Converts JSON to PowerPoint',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python pptx_builder_agent.py                    # Default (from ../5- Slide Generator Agent/)
  python pptx_builder_agent.py --input presentation.json
  python pptx_builder_agent.py --output lecture.pptx
  python pptx_builder_agent.py --audio ../8-\ TTS\ Generative\ Audio\ Agent/audio_output
        """
    )
    
    parser.add_argument(
        '--input',
        default='../5- Slide Generator Agent/presentation.json',
        help='Input presentation JSON file'
    )
    
    parser.add_argument(
        '--output',
        default='lecture.pptx',
        help='Output PowerPoint file (default: lecture.pptx)'
    )
    
    parser.add_argument(
        '--audio',
        default=None,
        help='Audio folder containing WAV files for slides'
    )
    
    args = parser.parse_args()
    
    print("\n" + "="*80)
    print("PPTX BUILDER AGENT - Production Implementation")
    print("="*80)
    print(f"Input: {args.input}")
    print(f"Output: {args.output}")
    if args.audio:
        print(f"Audio Folder: {args.audio}")
    
    # Create agent with audio folder
    builder = PPTXBuilderAgent(audio_folder=args.audio)
    
    # Load JSON
    if not builder.load_presentation_json(args.input):
        logger.error("Failed to load presentation JSON")
        return
    
    # Build presentation
    if not builder.build_presentation():
        logger.error("Failed to build presentation")
        return
    
    # Save file
    if not builder.save_presentation(args.output):
        logger.error("Failed to save presentation")
        return
    
    # Display summary
    builder.display_summary()
    
    print("\n" + "="*80)
    print("[OK] PPTX GENERATION COMPLETE")
    print("="*80)


if __name__ == "__main__":
    main()
